from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

import fitz
from docx import Document
from docx.enum.text import WD_BREAK
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "xlsx"}
MIME_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


def fail(code: str, message: str) -> None:
    print(json.dumps({"ok": False, "code": code, "message": message}, ensure_ascii=False))
    raise SystemExit(1)


def safe_leaf(value: str) -> str:
    leaf = Path(value.replace("\\", "/")).name
    if not leaf or leaf in {".", ".."} or "\x00" in leaf:
        fail("artifact_output_name_invalid", "The output file name is invalid.")
    return leaf[:180]


def output_extension(name: str) -> str:
    extension = Path(name).suffix.lower().lstrip(".")
    if extension not in SUPPORTED_EXTENSIONS:
        fail("artifact_format_unsupported", f"Unsupported artifact format: {extension or 'none'}")
    return extension


def read_request() -> dict[str, Any]:
    request_path = Path(os.environ.get("MY_MATE_ARTIFACT_REQUEST", "/input/request.json"))
    try:
        value = json.loads(request_path.read_text(encoding="utf-8"))
    except Exception:
        fail("artifact_request_invalid", "Artifact Worker could not read a valid request.")
    if not isinstance(value, dict):
        fail("artifact_request_invalid", "Artifact Worker request must be an object.")
    return value


def contains_cjk(value: str) -> bool:
    return bool(re.search(r"[\u3400-\u9fff]", value))


def register_document_fonts(content: str) -> tuple[str, str | None]:
    needs_cjk = contains_cjk(content)
    latin_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    if not Path(latin_path).is_file():
        fail("artifact_font_unavailable", "Artifact Worker could not load an embeddable document font.")
    try:
        pdfmetrics.registerFont(TTFont("ArtifactSans", latin_path))
    except Exception:
        fail("artifact_font_unavailable", "Artifact Worker could not load an embeddable document font.")
    if not needs_cjk:
        return "ArtifactSans", None
    cjk_path = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
    if not Path(cjk_path).is_file():
        fail("artifact_cjk_font_unavailable", "Artifact Worker could not load an embeddable CJK font.")
    try:
        pdfmetrics.registerFont(TTFont("ArtifactCJK", cjk_path))
    except Exception:
        fail("artifact_cjk_font_unavailable", "Artifact Worker could not load an embeddable CJK font.")
    return "ArtifactSans", "ArtifactCJK"


def format_pdf_text(value: str, cjk_font: str | None) -> str:
    escaped = value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    if not cjk_font:
        return escaped
    return re.sub(
        r"([\u2e80-\u9fff\uf900-\ufaff\u3040-\u30ff\uac00-\ud7af]+)",
        rf'<font name="{cjk_font}">\1</font>',
        escaped,
    )


def clean_inline_markdown(value: str) -> str:
    value = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", value)
    value = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", value)
    value = re.sub(r"`([^`]+)`", r"\1", value)
    value = re.sub(r"[*_~]+", "", value)
    return value.strip()


def content_lines(content: str) -> list[str]:
    return content.replace("\r\n", "\n").replace("\r", "\n").split("\n")


def create_pdf(content: str, output_path: Path, title: str) -> None:
    font, cjk_font = register_document_fonts(content)
    styles = getSampleStyleSheet()
    body = ParagraphStyle("ArtifactBody", parent=styles["BodyText"], fontName=font, fontSize=10.5, leading=16, spaceAfter=6)
    heading = {
        1: ParagraphStyle("ArtifactH1", parent=body, fontSize=22, leading=28, spaceBefore=8, spaceAfter=12),
        2: ParagraphStyle("ArtifactH2", parent=body, fontSize=16, leading=22, spaceBefore=12, spaceAfter=8),
        3: ParagraphStyle("ArtifactH3", parent=body, fontSize=13, leading=18, spaceBefore=10, spaceAfter=6),
    }
    code = ParagraphStyle("ArtifactCode", parent=body, fontName=font, fontSize=8.5, leading=12, backColor=colors.HexColor("#F3F4F6"), borderPadding=6)
    document = SimpleDocTemplate(
        str(output_path), pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm,
        topMargin=18 * mm, bottomMargin=18 * mm, title=title,
    )
    story: list[Any] = []
    in_code = False
    for raw in content_lines(content):
        stripped = raw.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if stripped in {"---", "***", "___"}:
            story.append(Spacer(1, 5 * mm))
            continue
        match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if match:
            level = min(3, len(match.group(1)))
            story.append(Paragraph(format_pdf_text(clean_inline_markdown(match.group(2)), cjk_font), heading[level]))
            continue
        if stripped == "\f":
            story.append(PageBreak())
            continue
        if not stripped:
            story.append(Spacer(1, 2.5 * mm))
            continue
        list_item = re.match(r"^([-*+]|\d+[.)])\s+(.*)$", stripped)
        text = clean_inline_markdown(list_item.group(2) if list_item else stripped)
        prefix = ""
        if list_item:
            marker = list_item.group(1)
            prefix = f"{marker if marker[0].isdigit() else '-'} "
        story.append(Paragraph(prefix + format_pdf_text(text, cjk_font), code if in_code else body))
    if not story:
        story.append(Paragraph(format_pdf_text(title or "Artifact", cjk_font), heading[1]))
    document.build(story)


def create_docx(content: str, output_path: Path, title: str) -> None:
    document = Document()
    document.core_properties.title = title
    document.core_properties.author = "My Mate Studio"
    in_code = False
    for raw in content_lines(content):
        stripped = raw.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if match:
            document.add_heading(clean_inline_markdown(match.group(2)), level=min(9, len(match.group(1))))
            continue
        bullet = re.match(r"^[-*+]\s+(.*)$", stripped)
        numbered = re.match(r"^\d+[.)]\s+(.*)$", stripped)
        if bullet:
            document.add_paragraph(clean_inline_markdown(bullet.group(1)), style="List Bullet")
        elif numbered:
            document.add_paragraph(clean_inline_markdown(numbered.group(1)), style="List Number")
        elif stripped == "\f":
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        elif not stripped:
            document.add_paragraph()
        else:
            paragraph = document.add_paragraph(clean_inline_markdown(stripped))
            if in_code:
                for run in paragraph.runs:
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
    document.save(output_path)


def markdown_sections(content: str, fallback_title: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    title = fallback_title
    body: list[str] = []
    for raw in content_lines(content):
        match = re.match(r"^#{1,2}\s+(.*)$", raw.strip())
        if match:
            if body or sections:
                sections.append((title, body))
            title = clean_inline_markdown(match.group(1)) or fallback_title
            body = []
        elif raw.strip():
            body.append(clean_inline_markdown(raw.strip()))
    if body or not sections:
        sections.append((title, body))
    return sections[:80]


def create_pptx(content: str, output_path: Path, title: str) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    sections = markdown_sections(content, title or "Artifact")
    for index, (section_title, lines) in enumerate(sections):
        layout = presentation.slide_layouts[0 if index == 0 else 1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = section_title
        if index == 0:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = next((line for line in lines if line), "Generated by My Mate Studio")
            continue
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for line_index, line in enumerate(lines[:12]):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.text = re.sub(r"^(?:[-*+]\s+|\d+[.)]\s+)", "", line)
            paragraph.level = 0
            paragraph.font.size = Pt(20)
    presentation.save(output_path)


def parse_table_content(content: str) -> tuple[list[str], list[list[Any]]]:
    try:
        value = json.loads(content)
        if isinstance(value, dict) and isinstance(value.get("columns"), list) and isinstance(value.get("rows"), list):
            return [str(item) for item in value["columns"]], [list(row) for row in value["rows"] if isinstance(row, list)]
    except Exception:
        pass
    rows = [line.split("\t") for line in content_lines(content) if line.strip()]
    if rows:
        width = max(len(row) for row in rows)
        columns = rows[0] if len(rows) > 1 else [f"Column {index + 1}" for index in range(width)]
        return columns, rows[1:] if len(rows) > 1 else rows
    return ["Content"], [[content]]


def create_xlsx(content: str, output_path: Path, title: str) -> None:
    columns, rows = parse_table_content(content)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (title or "Sheet1")[:31]
    sheet.append(columns)
    for row in rows:
        sheet.append(row[: len(columns)] + [None] * max(0, len(columns) - len(row)))
    sheet.freeze_panes = "A2"
    sheet.auto_filter.ref = sheet.dimensions
    for column in sheet.columns:
        letter = column[0].column_letter
        sheet.column_dimensions[letter].width = min(48, max(12, max(len(str(cell.value or "")) for cell in column) + 2))
    workbook.save(output_path)


def libreoffice_convert(source: Path, extension: str, output_dir: Path) -> Path:
    profile = tempfile.mkdtemp(prefix="lo-profile-", dir="/tmp")
    try:
        completed = subprocess.run(
            [
                "libreoffice", "--headless", "--nologo", "--nodefault", "--nolockcheck", "--nofirststartwizard",
                f"-env:UserInstallation=file://{profile}", "--convert-to", extension, "--outdir", str(output_dir), str(source),
            ],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=90, check=False,
        )
        expected = output_dir / f"{source.stem}.{extension}"
        if completed.returncode != 0 or not expected.is_file():
            fail("artifact_conversion_failed", "LibreOffice did not produce the requested artifact.")
        return expected
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def extract_text(path: Path, extension: str) -> str:
    if extension == "pdf":
        document = fitz.open(path)
        try:
            return "\n\n".join(page.get_text("text").strip() for page in document if page.get_text("text").strip())
        finally:
            document.close()
    if extension == "docx":
        document = Document(path)
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            lines.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        return "\n".join(lines)
    if extension == "pptx":
        presentation = Presentation(path)
        return "\n\n".join(
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            for slide in presentation.slides
        ).strip()
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                lines.append("\t".join("" if cell is None else str(cell) for cell in row))
        return "\n".join(lines).strip()
    finally:
        workbook.close()


def validate_output(path: Path, extension: str) -> dict[str, Any]:
    if not path.is_file() or path.stat().st_size <= 0:
        fail("artifact_output_missing", "Artifact Worker did not create a non-empty output file.")
    if extension == "pdf":
        reader = PdfReader(str(path))
        if not reader.pages:
            fail("artifact_validation_failed", "The generated PDF has no pages.")
        details = {"page_count": len(reader.pages)}
    elif extension == "docx":
        document = Document(path)
        details = {"paragraph_count": len(document.paragraphs), "table_count": len(document.tables)}
    elif extension == "pptx":
        presentation = Presentation(path)
        if not presentation.slides:
            fail("artifact_validation_failed", "The generated presentation has no slides.")
        details = {"slide_count": len(presentation.slides)}
    else:
        workbook = load_workbook(path, read_only=True, data_only=False)
        try:
            if not workbook.sheetnames:
                fail("artifact_validation_failed", "The generated workbook has no worksheets.")
            details = {"sheet_count": len(workbook.sheetnames), "sheet_names": workbook.sheetnames[:20]}
        finally:
            workbook.close()
    return details


def validate_pdf_rendering(
    path: Path,
    source_content: str,
    extracted_text: str,
    require_embedded_font: bool,
) -> dict[str, Any]:
    document = fitz.open(path)
    try:
        if not document.page_count:
            fail("artifact_validation_failed", "The generated PDF has no pages.")

        source_cjk_characters = len(re.findall(r"[\u3400-\u9fff]", source_content))
        extracted_cjk_characters = len(re.findall(r"[\u3400-\u9fff]", extracted_text))
        if source_cjk_characters and extracted_cjk_characters < min(4, source_cjk_characters):
            fail(
                "artifact_cjk_text_validation_failed",
                "The generated PDF did not preserve extractable CJK text.",
            )

        font_xrefs: set[int] = set()
        embedded_font_xrefs: set[int] = set()
        for page in document:
            for font in page.get_fonts(full=True):
                xref = int(font[0])
                if xref <= 0 or xref in font_xrefs:
                    continue
                font_xrefs.add(xref)
                try:
                    _name, _extension, _font_type, font_bytes = document.extract_font(xref)
                    if font_bytes:
                        embedded_font_xrefs.add(xref)
                except Exception:
                    continue
        if require_embedded_font and not embedded_font_xrefs:
            fail("artifact_font_embedding_failed", "The generated PDF does not contain an embedded font.")

        page = document.load_page(0)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), colorspace=fitz.csGRAY, alpha=False)
        samples = pixmap.samples
        ink_pixels = sum(1 for value in samples if value < 245)
        total_pixels = max(1, pixmap.width * pixmap.height)
        minimum_ink_pixels = max(128, int(total_pixels * 0.0002))
        if ink_pixels < minimum_ink_pixels:
            fail("artifact_pdf_blank_page", "The generated PDF rendered as a blank or nearly blank first page.")
        return {
            "source_cjk_characters": source_cjk_characters,
            "extracted_cjk_characters": extracted_cjk_characters,
            "font_count": len(font_xrefs),
            "embedded_font_count": len(embedded_font_xrefs),
            "rendered_page": 1,
            "render_width": pixmap.width,
            "render_height": pixmap.height,
            "render_ink_pixels": ink_pixels,
            "render_ink_ratio": round(ink_pixels / total_pixels, 8),
        }
    finally:
        document.close()


def make_preview(output_path: Path, extension: str, output_dir: Path) -> Path | None:
    if extension == "pdf":
        return output_path
    preview_dir = output_dir / "preview"
    preview_dir.mkdir(parents=True, exist_ok=True)
    converted = libreoffice_convert(output_path, "pdf", preview_dir)
    preview = output_dir / "preview.pdf"
    converted.replace(preview)
    return preview


def main() -> None:
    request = read_request()
    output_dir = Path("/output")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_name = safe_leaf(str(request.get("output_name") or ""))
    extension = output_extension(output_name)
    output_path = output_dir / output_name
    content = str(request.get("content") or "")
    title = str(request.get("title") or Path(output_name).stem or "Artifact")[:200]
    source_value = str(request.get("source_file") or "").strip()
    source_path = Path("/input") / safe_leaf(source_value) if source_value else None

    converted_directly = False
    if source_path and source_path.is_file() and bool(request.get("prefer_source_conversion")):
        source_extension = source_path.suffix.lower().lstrip(".")
        if source_extension == extension:
            shutil.copyfile(source_path, output_path)
            converted_directly = True
        elif source_extension in SUPPORTED_EXTENSIONS:
            converted = libreoffice_convert(source_path, extension, output_dir)
            if converted != output_path:
                converted.replace(output_path)
            converted_directly = True

    if not converted_directly:
        if not content.strip() and source_path and source_path.is_file():
            source_extension = source_path.suffix.lower().lstrip(".")
            if source_extension in SUPPORTED_EXTENSIONS:
                content = extract_text(source_path, source_extension)
        if not content.strip():
            fail("artifact_content_missing", "No content was available for artifact generation.")
        if extension == "pdf":
            create_pdf(content, output_path, title)
        elif extension == "docx":
            create_docx(content, output_path, title)
        elif extension == "pptx":
            create_pptx(content, output_path, title)
        else:
            create_xlsx(content, output_path, title)

    validation = validate_output(output_path, extension)
    extracted_text = extract_text(output_path, extension).replace("\x00", "")
    extractable_text = extracted_text.strip()
    if not converted_directly and content.strip() and not extractable_text:
        fail("artifact_text_validation_failed", "The generated artifact does not contain extractable text.")
    validation["extracted_text_characters"] = len(extractable_text)
    if extension == "pdf":
        validation.update(validate_pdf_rendering(
            output_path,
            content,
            extracted_text,
            require_embedded_font=not converted_directly,
        ))
    preview_path = make_preview(output_path, extension, output_dir)
    digest = hashlib.sha256(output_path.read_bytes()).hexdigest()
    manifest = {
        "ok": True,
        "schema_version": 1,
        "worker_version": "0.1.0",
        "output_file": output_path.name,
        "mime_type": MIME_TYPES[extension],
        "size_bytes": output_path.stat().st_size,
        "sha256": digest,
        "preview_file": preview_path.name if preview_path else None,
        "extracted_text_file": "extracted.txt",
        "validation": validation,
    }
    (output_dir / "extracted.txt").write_text(extracted_text[:2_000_000], encoding="utf-8")
    (output_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False), encoding="utf-8")
    print(json.dumps(manifest, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except SystemExit:
        raise
    except Exception:
        fail("artifact_worker_failed", "Artifact Worker failed without exposing private content.")
